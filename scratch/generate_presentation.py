import os
import sys
from pptx import Presentation

# Shape-based mapping for robust, layout-preserving text replacement
# Key format: (slide_index_1_based, shape_index_0_based) -> new_text
SHAPE_MAPPING = {
    # ==================== SLIDE 1 ====================
    (1, 4): "AI-NATIVE LEARNING COMPANION",
    (1, 5): "StudyPal",
    (1, 6): "Socratic tutoring, smart planning, active recall, and auditory learning — all in one accessible platform.",
    (1, 8): "Next.js  ·  FastAPI  ·  PostgreSQL  ·  qwen2.5:7b  ·  gemma4:e2b  ·  nomic-embed-text  ·  duckduckgo  ·  Vocal Bridge",

    # ==================== SLIDE 2 ====================
    (2, 1): "THE CHALLENGE",
    (2, 2): "Traditional learning\nis rigid and expensive.",
    (2, 5): "One-size-fits-all classroom pacing triggers executive dysfunction and mental fatigue",
    (2, 8): "High financial and scheduling barriers to access specialized TAs or private tutors",
    (2, 11): "Learning anxiety prevents students from asking repeated questions or admitting confusion",
    (2, 14): "Static textbooks and worksheets fail to adapt to diverse visual or auditory processing styles",
    (2, 17): "Lack of active engagement and distraction blocks deep, productive focus",
    (2, 19): "Without StudyPal",
    (2, 21): "how do I solve this?",
    (2, 23): "just memorize it 🙄",
    (2, 25): "I need a tutor now",
    (2, 27): "exam tomorrow?",
    (2, 29): "I'm completely lost",
    (2, 31): "where do I start?",
    (2, 33): "too much text, closing it",

    # ==================== SLIDE 3 ====================
    (3, 1): "PRODUCT FEATURES",
    (3, 2): "An all-in-one suite designed to accommodate every learning style.",
    (3, 6): "Configurable Tutor Bot",
    (3, 7): "Highly personalized companion configured with custom personality (Soul), user profile contexts, tools, agent instructions, and background heartbeat tasks.",
    (3, 11): "Socratic Chat & Agent",
    (3, 12): "Conversational environment running deep capabilities: Socratic chat, deepsolve, quiz generation, deep research, and math animator.",
    (3, 16): "Lightweight Agent Tools",
    (3, 17): "On-demand execution tools called dynamically by the AI, including brainstorm, rag, web search, code execution, reason, and arXiv search.",
    (3, 21): "Vocal Bridge (Voice)",
    (3, 22): "Seamless real-time conversational audio agent for hands-free and auditory-first learning (perfect for dyslexia accommodations).",
    (3, 26): "Knowledge Vault (RAG)",
    (3, 27): "Upload and index textbooks, lecture slides, and PDFs to search and retrieve context instantly through custom knowledge bases.",
    (3, 31): "Episodic Memory",
    (3, 32): "A secure local memory module tracking user explanation preferences, custom rules, and study profiles to maintain context across sessions.",

    # ==================== SLIDE 4 ====================
    (4, 1): "TECHNICAL ARCHITECTURE",
    (4, 2): "How StudyPal is built.",
    (4, 5): "Web Workspace",
    (4, 6): "Next.js  ·  Tailwind CSS  ·  TypeScript",
    (4, 10): "Responsive, modern React application. Implements workspaces for tools like Co-Writer, Focus Mode, Decks, and Interactive Boards.",
    (4, 14): "Agent Engine",
    (4, 15): "Python  ·  FastAPI  ·  Orchestrator",
    (4, 19): "Runs the dual-layer agent system. Level 1 tools (RAG, Web search) and Level 2 Capabilities (Deep Solve, Deep Question) route dynamically.",
    (4, 23): "AI Models & Voice",
    (4, 24): "qwen2.5:7b  ·  gemma4:e2b  ·  Vocal Bridge",
    (4, 28): "Configured with local models for deep reasoning and math animations, alongside Vocal Bridge APIs for direct voice interaction.",
    (4, 30): "SYSTEM STACK",
    (4, 31): "Host",
    (4, 32): "Vercel & Railway",
    (4, 34): "Frontend",
    (4, 35): "Next.js Workspace",
    (4, 37): "Backend",
    (4, 38): "FastAPI / Python",
    (4, 40): "Voice",
    (4, 41): "Vocal Bridge Agent",
    (4, 43): "Memory",
    (4, 44): "Vector & Relational DB",
    (4, 46): "Orchestration",
    (4, 47): "Dynamic Router",

    # ==================== SLIDE 5 ====================
    (5, 1): "COGNITIVE IMPACT",
    (5, 2): "Empowering independent, self-paced learning.",
    (5, 3): "THE 10 WORKSPACE TOOLS & METHODS",
    (5, 7): "Visual Mastery: Whiteboard (Infinite sketch canvas with GeoGebra curve graphing), interactive Mindmap generator, and automatic Slide Deck creator.",
    (5, 11): "Audio & Productivity: Podcast generator (converts text/notes to script into 2-person discussions), Adaptive Study Planner, Focus Mode, and Co-Writer (blank-page aid).",
    (5, 15): "Active Recall & Exam Prep: Guided Learning (provides structured learning through rich, interactive HTML pages), Spaced-Repetition Flash Cards, and Exam Simulator.",
    (5, 19): "Accessible Tutoring: Removes scheduling stress, eliminates textbook reading fatigue, and mitigates test anxiety through highly engaging Socratic practice.",
    (5, 21): "STUDENT ENGAGEMENT IMPACT",
    (5, 24): "24",
    (5, 25): "/7",
    (5, 26): "unlimited access to patient, empathetic Socratic tutoring",
    (5, 29): "0",
    (5, 30): "dollars",
    (5, 31): "spent on expensive, scheduling-restricted private tutors",
    (5, 34): "2x",
    (5, 35): "increase in comprehension and long-term active recall",
    (5, 38): "100",
    (5, 39): "%",
    (5, 40): "learning autonomy tailored to neurodivergent pacing"
}

def replace_text_preserve_formatting(shape, replacement_text):
    if not shape.has_text_frame:
        return False
        
    tf = shape.text_frame
    
    # We clear subsequent runs/paragraphs to maintain one styled block.
    # To keep formatting, we write to the first run of the first paragraph, and clear others.
    if tf.paragraphs:
        p = tf.paragraphs[0]
        # If paragraph has runs, use the first run's font/color formatting
        if p.runs:
            p.runs[0].text = replacement_text
            # Clear text from all subsequent runs in this paragraph
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = replacement_text
            
        # Clear text from all other paragraphs in the text frame
        for other_p in tf.paragraphs[1:]:
            for r in other_p.runs:
                r.text = ""
            other_p.text = ""
    else:
        tf.text = replacement_text
    return True

def build_presentation(input_path, output_path):
    print(f"Opening template: {input_path}")
    if not os.path.exists(input_path):
        print(f"Error: template file [{input_path}] does not exist!")
        return False
        
    prs = Presentation(input_path)
    print(f"Loaded template presentation containing {len(prs.slides)} slides.")
    
    success_count = 0
    
    for (slide_num, shape_idx), replacement_text in SHAPE_MAPPING.items():
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            print(f"Warning: slide index {slide_idx} out of bounds.")
            continue
            
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            print(f"Warning: shape index {shape_idx} out of bounds on slide {slide_num}.")
            continue
            
        shape = slide.shapes[shape_idx]
        if not shape.has_text_frame:
            print(f"Warning: shape index {shape_idx} on slide {slide_num} does not have a text frame.")
            continue
            
        old_text = " ".join(p.text for p in shape.text_frame.paragraphs if p.text)
        print(f"Slide {slide_num}, Shape {shape_idx}: Replacing [{old_text[:30]}...] -> [{replacement_text[:30]}...]")
        
        if replace_text_preserve_formatting(shape, replacement_text):
            success_count += 1
            
    print(f"\nProcessing complete. Replaced {success_count} fields successfully.")
    print(f"Saving presentation to: {output_path}")
    prs.save(output_path)
    print("Success! File saved successfully.")
    return True

if __name__ == "__main__":
    build_presentation("template.pptx", "studypal_presentation.pptx")
