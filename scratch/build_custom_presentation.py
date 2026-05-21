import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Theme Colors (Reflecting StudyPal's Official Light Theme)
BG_COLOR = RGBColor(248, 250, 252)       # #F8FAFC Light Canvas Background
CARD_COLOR = RGBColor(255, 255, 255)     # #FFFFFF Card Surface
BORDER_COLOR = RGBColor(226, 232, 240)   # #E2E8F0 Subtle Border Outline
TEXT_DARK = RGBColor(15, 23, 42)         # #0F172A Slate 900 Title/Main Text
TEXT_GRAY = RGBColor(100, 116, 139)      # #64748B Slate 500 Secondary/Muted Text
ACCENT_PURPLE = RGBColor(124, 58, 237)   # #7C3AED Rich Indigo/Violet
ACCENT_BLUE = RGBColor(37, 99, 235)      # #2563EB Vibrant Accent Blue
ACCENT_GREEN = RGBColor(5, 150, 105)     # #059669 Emerald Green for Success/Metrics

def add_background(slide, prs):
    # Create a full slide rectangle for the background
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_COLOR
    shape.line.fill.background() # No border for the background canvas

def add_card(slide, left, top, width, height, title, content, title_color=ACCENT_PURPLE, align_center=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    # Add light border to make cards pop against the off-white background
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(1.5)
    
    # Unified text box to prevent overlapping/clipping and support natural wrapping
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.05), width - Inches(0.3), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
    
    # Determine sizing based on card height
    if height < Inches(1.0):
        title_size = Pt(13)
        content_size = Pt(11)
        space_before = Pt(2)
    else:
        title_size = Pt(16)
        content_size = Pt(12)
        space_before = Pt(6)
        
    # Title Paragraph
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = title_size
    p.font.color.rgb = title_color
    if align_center:
        p.alignment = PP_ALIGN.CENTER
        
    # Content Paragraph
    if content:
        p_c = tf.add_paragraph()
        p_c.text = content
        p_c.font.size = content_size
        p_c.font.color.rgb = TEXT_DARK  # Standard dark slate text
        p_c.space_before = space_before
        if align_center:
            p_c.alignment = PP_ALIGN.CENTER

def add_metric_card(slide, left, top, width, height, metric, description, metric_color=ACCENT_PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(1.5)
    
    # Metric (Huge)
    txBox = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = metric
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = metric_color
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    if description:
        txBox_c = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.0), width - Inches(0.2), Inches(0.8))
        tf_c = txBox_c.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.add_paragraph()
        p_c.text = description
        p_c.font.size = Pt(12)
        p_c.font.color.rgb = TEXT_DARK
        p_c.alignment = PP_ALIGN.CENTER

def add_header(slide, title, subtitle):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = TEXT_DARK
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_GRAY

def add_image_safely(slide, image_path, left, top, width, height):
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"Embedded image: {image_path}")
        except Exception as e:
            print(f"Failed to embed image {image_path}: {e}")
            # Add a placeholder card if image load fails
            add_card(slide, left, top, width, height, "[Image Load Failed]", f"Path: {image_path}\nError: {e}", TEXT_GRAY)
    else:
        print(f"Image not found: {image_path}")
        # Add a visual warning card if the asset is missing
        add_card(slide, left, top, width, height, "[Missing Image Asset]", f"Please place the file at:\n{image_path}", TEXT_GRAY)

def build_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # --- Slide 1: Title Slide ---
    s1 = prs.slides.add_slide(blank_layout)
    add_background(s1, prs)
    
    txBox = s1.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.33), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "StudyPal"
    p.font.bold = True
    p.font.size = Pt(72)
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "The AI-Native Learning Operating System for Neurodivergent Minds."
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_PURPLE
    p2.alignment = PP_ALIGN.CENTER
    
    txBox_footer = s1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
    tf_f = txBox_footer.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "Next.js · FastAPI · PostgreSQL · qwen2.5:7b · gemma4:e2b · nomic-embed-text · duckduckgo · Vocal Bridge"
    p_f.font.size = Pt(14)
    p_f.font.color.rgb = TEXT_GRAY
    p_f.alignment = PP_ALIGN.CENTER

    # --- Slide 2: The Accessibility Gap ---
    s2 = prs.slides.add_slide(blank_layout)
    add_background(s2, prs)
    add_header(s2, "THE ACCESSIBILITY GAP IN MODERN LEARNING", "Standardized education creates massive barriers for ADHD, dyslexia, and neurodivergent learners.")
    
    add_card(s2, Inches(1.0), Inches(2.0), Inches(5.2), Inches(2.1), "1. The Pacing & Overload Trap", "Standard classrooms move too fast, triggering executive dysfunction and cognitive fatigue. Static textbooks fail visual or auditory-first processing styles.", ACCENT_PURPLE)
    add_card(s2, Inches(7.13), Inches(2.0), Inches(5.2), Inches(2.1), "2. The Dependency & Financial Barrier", "Accessing specialized support like private tutors or TAs is highly expensive, scheduling-restricted, and hard to coordinate.", ACCENT_BLUE)
    add_card(s2, Inches(1.0), Inches(4.6), Inches(5.2), Inches(2.1), "3. The Learning Anxiety Loop", "Fear of peer judgment or admitting confusion stops students from asking questions, causing them to fall further behind.", ACCENT_BLUE)
    add_card(s2, Inches(7.13), Inches(4.6), Inches(5.2), Inches(2.1), "4. Lack of Active Engagement", "Passive reading or listening fails to keep learners with ADHD engaged, leading to rapid distraction and loss of interest.", ACCENT_PURPLE)

    # --- Slide 3: The Solution (with Socratic Chat Workspace Screenshot) ---
    s3 = prs.slides.add_slide(blank_layout)
    add_background(s3, prs)
    add_header(s3, "THE SOLUTION: STUDYPAL (DEEPTUTOR)", "An agent-native, personalized cognitive companion empowering independent mastery.")
    
    # Left column - Core Pillars
    add_card(s3, Inches(0.5), Inches(2.0), Inches(4.5), Inches(1.4), "Infinite Patience & Safe Space", "Allows students to ask questions 10 different ways at 2:00 AM with zero frustration or peer judgment.", ACCENT_PURPLE)
    add_card(s3, Inches(0.5), Inches(3.65), Inches(4.5), Inches(1.4), "Visual Scaffolding (Manim & GeoGebra)", "Transforms abstract formulas into beautiful dynamic animations (Manim) and interactive graph canvas (GeoGebra).", ACCENT_BLUE)
    add_card(s3, Inches(0.5), Inches(5.3), Inches(4.5), Inches(1.4), "Socratic Guides (deep_solve & memory)", "Breaks down massive problems into bite-sized cognitive stages. Recalls preferences and progress.", ACCENT_GREEN)
    
    # Middle column - Workspace description
    add_card(s3, Inches(5.3), Inches(2.0), Inches(3.6), Inches(4.7), "Copilot & Voice Integration", 
             "• Socratic Conversational Bot: Guides students toward solutions instead of giving answers.\n\n"
             "• Auditory Voice Integration: Hands-free conversation. Perfect for dyslexia or auditory-first study habits.\n\n"
             "• Hyper-Personalized Knowledge Vault: Instantly uploads textbooks, slides, and class notes to construct localized, authoritative expert brains.", ACCENT_PURPLE)
             
    # Right column - Live Screenshot
    add_image_safely(s3, "assets/figs/dt-chat.png", Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.7))

    # --- Slide 4: Specialized Workspace Tools (with Co-Writer and Knowledge Vault Screenshots) ---
    s4 = prs.slides.add_slide(blank_layout)
    add_background(s4, prs)
    add_header(s4, "THE INTERACTIVE STUDENT WORKSPACE", "Specialized productivity, visualization, and active recall suites built into a single workspace.")
    
    # Left side: Visual & Auditory, and Productivity
    add_card(s4, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.2), "Visual & Audio Mastery", 
             "• Whiteboard: Infinite digital canvas rendering interactive GeoGebra curves.\n"
             "• Mindmap: Parses text into visual semantic node graphs.\n"
             "• Slide Deck: Dynamically generates structured slides from study notes.\n"
             "• Podcast: Conversational text-to-podcast script and discussion engine.", ACCENT_PURPLE)
             
    add_card(s4, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.2), "Productivity & Focus", 
             "• Study Planner: Adaptive organizer aligning with personal energy levels.\n"
             "• Focus Mode: Distraction-free space blocking cognitive overload.\n"
             "• Co-Writer: Cooperative drafting tool for overcoming blank-page anxiety.", ACCENT_BLUE)
             
    # Right side top: Active Recall
    add_card(s4, Inches(6.8), Inches(2.0), Inches(6.0), Inches(2.2), "Active Recall & Exam Preparation", 
             "• Guided Learning: Socratic curriculum paths rendered as rich interactive HTML pages.\n"
             "• Flash Cards: Spaced-repetition card decks generated directly from notes.\n"
             "• Exam Simulator: Dynamic mock exams with automated grading, diagnostic reports, and error reviews.", ACCENT_GREEN)
             
    # Right side bottom: Screenshots of Co-Writer and Knowledge Vault side-by-side
    add_image_safely(s4, "assets/figs/dt-cowriter.png", Inches(6.8), Inches(4.6), Inches(2.9), Inches(2.2))
    add_image_safely(s4, "assets/figs/dt-knowledge.png", Inches(9.9), Inches(4.6), Inches(2.9), Inches(2.2))

    # --- Slide 5: Backend Architecture (with deeptutor-architecture.png) ---
    s5 = prs.slides.add_slide(blank_layout)
    add_background(s5, prs)
    add_header(s5, "ARCHITECTURE: CONVERSATIONAL ENGINE", "An asynchronous orchestrated backend managing real-time Socratic interactions.")
    
    # Left side: text explanation
    add_card(s5, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.8), "Orchestrated Multi-Agent Execution",
             "• Unified ChatOrchestrator: unified entry point routing WebSocket requests to standard tools or multi-step capabilities.\n\n"
             "• Level 1 Tools (The Muscle): Stateless, single-purpose LLM-called utilities executing immediate actions (e.g. RAG lookup, web_search, code_execution).\n\n"
             "• Level 2 Capabilities (The Brain): Stateful multi-step agent pipelines taking over the conversation flow (e.g., deep_solve path reasoning, deep_question validation).\n\n"
             "• StreamBus Event Hub: Async event fan-out streaming intermediate thoughts, tool invocations, and answers to the websocket interface.", ACCENT_PURPLE)
             
    # Right side: Architecture Diagram
    add_image_safely(s5, "assets/figs/deeptutor-architecture.png", Inches(6.7), Inches(2.0), Inches(6.1), Inches(4.8))

    # --- Slide 6: Frontend Architecture (Copilot Loop) ---
    s6 = prs.slides.add_slide(blank_layout)
    add_background(s6, prs)
    add_header(s6, "ARCHITECTURE: COPILOT WORKSPACE LOOP", "Deeply integrating AI directly into React component states for active workspace operations.")
    
    add_card(s6, Inches(0.5), Inches(2.3), Inches(3.8), Inches(3.4), "1. Generative UI (CopilotKit)", "Wraps the Next.js frontend, managing real-time chat state and LLM connectivity (Gemma/Qwen/local models) across active tabs.", ACCENT_PURPLE)
    add_card(s6, Inches(4.76), Inches(2.3), Inches(3.8), Inches(3.4), "2. Readable State (useCopilotReadable)", "Exposes active workspace data (e.g., current focus timer status, whiteboard graph definitions, student schedule) directly to the AI context.", ACCENT_BLUE)
    add_card(s6, Inches(9.03), Inches(2.3), Inches(3.8), Inches(3.4), "3. Executable Actions (useCopilotAction)", "Enables the AI agent to programmatically execute UI tasks (e.g., automatically scheduling focus blocks or generating slides on the canvas).", ACCENT_GREEN)
    
    # Adding a visual layout connector text box
    tx = s6.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = "Feedback Loop: Read state -> Reason (Orchestrator) -> Execute Action -> Render updated state."
    p.font.italic = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 7: Cognitive Impact & Infrastructure (with Memory dashboard Screenshot) ---
    s7 = prs.slides.add_slide(blank_layout)
    add_background(s7, prs)
    add_header(s7, "COGNITIVE IMPACT & DEPLOYMENT METRICS", "Quantified student outcomes and architectural efficiency.")
    
    # Metric cards taking less height
    add_metric_card(s7, Inches(0.5), Inches(2.0), Inches(2.8), Inches(1.9), "24/7", "Unlimited access to patient, empathetic Socratic tutoring.", ACCENT_PURPLE)
    add_metric_card(s7, Inches(3.7), Inches(2.0), Inches(2.8), Inches(1.9), "$0", "Token fee. Local model compilation eliminates Cloud API expenses.", ACCENT_BLUE)
    add_metric_card(s7, Inches(6.9), Inches(2.0), Inches(2.8), Inches(1.9), "2x+", "Increase in comprehension and long-term recall.", ACCENT_GREEN)
    add_metric_card(s7, Inches(10.1), Inches(2.0), Inches(2.7), Inches(1.9), "100%", "Learning privacy. Files and memories are kept locally.", ACCENT_PURPLE)
    
    # Wide screenshot of memory/progress dashboard at the bottom
    add_image_safely(s7, "assets/figs/dt-memory.png", Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.3))
    
    tx = s7.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "Privacy-First Learning OS: Runs entirely on consumer hardware (M-series Mac, local CUDA desktops)."
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    output_path = "studypal_custom_deck.pptx"
    prs.save(output_path)
    print(f"Successfully generated light theme custom presentation with screenshots: {output_path}")

if __name__ == "__main__":
    build_presentation()
