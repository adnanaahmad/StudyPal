import sys
from pptx import Presentation

def inspect_pptx(path):
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        return

    print(f"Number of slides: {len(prs.slides)}")
    print("\n--- SLIDES IN TEMPLATE ---")
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}: Layout = {slide.slide_layout.name}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(f"  Shape Type: {shape.name}, Text: {shape.text[:100]}")
            else:
                print(f"  Shape Type: {shape.name} (No text frame)")

    print("\n--- SLIDE MASTER LAYOUTS ---")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: Name = {layout.name}")
        for shape in layout.shapes:
            if shape.is_placeholder:
                print(f"  Placeholder index {shape.placeholder_format.idx}: name={shape.name}, type={shape.placeholder_format.type}")

if __name__ == "__main__":
    inspect_pptx("template.pptx")
